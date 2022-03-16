"""
Given a directory of powerpoint files, search the files for occurances of a
given regex.
"""
import csv
import re
from pathlib import Path

import pptx


INPUT_DIR = Path("infiles")
OUTPUT_CSV = Path(Path.home(), "Desktop", "output.csv")
REG = re.compile(
    r"http(s)?://musiclab.chromeexperiments.com/Song-Maker/song/(\d){16}",
)

matches = []
for f in Path.iterdir(INPUT_DIR):
    name = f.name.split(" - ")[0].strip()
    pres = pptx.Presentation(f)
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.replace("\n", "")
                if mo := re.search(REG, text):
                    matches.append((name, mo[0]))
            if hasattr(shape, "hyperlink"):
                raise ValueError(
                    "This script does not handle hyperlinks, but a hyperlink "
                    f"was found in presentation {f}. Add code here to search "
                    "the hyperlink for your regex pattern"
                )

with open(OUTPUT_CSV, "w") as csvf:
    wr = csv.writer(csvf)
    wr.writerow(["Name", "Link"])
    wr.writerows(matches)
